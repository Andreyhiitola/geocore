#!/usr/bin/env python3
"""Absolute-position pptx -> self-contained HTML converter.

Renders every shape at its own pptx coordinates (scaled to px) instead of
guessing slide semantics. Good for a fast visual diff/check after a partner
edits the .pptx by hand; not a substitute for the hand-built brand HTML.

Usage: python3 pptx_to_html.py input.pptx output.html
"""
import base64
import html
import io
import sys

import olefile
import openpyxl
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu

EMU_PER_PX = 914400 / 96  # 96 CSS px per inch


def px(emu):
    if emu is None:
        return 0
    return round(emu / EMU_PER_PX, 2)


def color_hex(color_format):
    try:
        if color_format.type is None:
            return None
        return '#' + str(color_format.rgb)
    except Exception:
        return None


def run_style(run):
    f = run.font
    parts = []
    if f.size:
        parts.append(f'font-size:{round(f.size.pt, 1)}px')
    if f.bold:
        parts.append('font-weight:700')
    if f.italic:
        parts.append('font-style:italic')
    if f.name:
        parts.append(f'font-family:"{f.name}",sans-serif')
    c = color_hex(f.color)
    if c:
        parts.append(f'color:{c}')
    return ';'.join(parts)


def render_text_frame(tf):
    out = []
    for p in tf.paragraphs:
        align = {1: 'center', 2: 'right', 3: 'justify'}.get(p.alignment, 'left') if p.alignment else 'left'
        runs_html = ''.join(
            f'<span style="{html.escape(run_style(r), quote=True)}">{html.escape(r.text)}</span>' for r in p.runs
        ) or '&nbsp;'
        out.append(f'<div style="text-align:{align};min-height:1em;">{runs_html}</div>')
    return ''.join(out)


def img_data_uri(image):
    b64 = base64.b64encode(image.blob).decode()
    return f'data:image/{image.ext};base64,{b64}'


def extract_ole_table_html(shape):
    """Best-effort: pull an embedded Excel OLE object out and render as a table."""
    try:
        ns = {
            'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
            'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
        }
        ole_obj = shape._element.find('.//p:oleObj', ns)
        rid = ole_obj.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
        data = shape.part.rels[rid].target_part.blob
    except Exception:
        return None
    try:
        ole = olefile.OleFileIO(io.BytesIO(data))
        pkg = ole.openstream('Package').read()
        wb = openpyxl.load_workbook(io.BytesIO(pkg), data_only=True)
        ws = wb.worksheets[0]
        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            return None
        head, *body = rows
        thead = ''.join(f'<th>{html.escape(str(c or ""))}</th>' for c in head)
        trs = []
        for r in body:
            if not any(c is not None for c in r):
                continue
            tds = ''.join(f'<td>{html.escape(str(c or ""))}</td>' for c in r)
            trs.append(f'<tr>{tds}</tr>')
        return f'<table class="ole-table"><thead><tr>{thead}</tr></thead><tbody>{"".join(trs)}</tbody></table>'
    except Exception:
        return None


def render_shape(shape):
    l, t, w, h = px(shape.left), px(shape.top), px(shape.width), px(shape.height)
    base_style = f'position:absolute;left:{l}px;top:{t}px;width:{w}px;height:{h}px;'

    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        try:
            src = img_data_uri(shape.image)
            return f'<img class="pshape" style="{base_style}object-fit:contain;" src="{src}">'
        except Exception:
            return ''

    if shape.has_table:
        rows_html = []
        for row in shape.table.rows:
            cells = ''.join(f'<td>{html.escape(c.text)}</td>' for c in row.cells)
            rows_html.append(f'<tr>{cells}</tr>')
        return f'<table class="pshape ptable" style="{base_style}"><tbody>{"".join(rows_html)}</tbody></table>'

    if shape.shape_type == MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT:
        table_html = extract_ole_table_html(shape)
        if table_html:
            return f'<div class="pshape ole-wrap" style="{base_style}overflow:auto;">{table_html}</div>'
        return f'<div class="pshape" style="{base_style}background:#eee;"></div>'

    fill = None
    try:
        if shape.fill.type is not None:
            fill = color_hex(shape.fill.fore_color)
    except Exception:
        pass
    style = base_style + (f'background:{fill};' if fill else '')

    if shape.has_text_frame and shape.text_frame.text.strip():
        return f'<div class="pshape ptext" style="{style}">{render_text_frame(shape.text_frame)}</div>'
    if fill:
        return f'<div class="pshape" style="{style}"></div>'
    return ''


def convert(pptx_path, out_path):
    prs = Presentation(pptx_path)
    sw, sh = px(prs.slide_width), px(prs.slide_height)

    slides_html = []
    for i, slide in enumerate(prs.slides, 1):
        shapes_html = []
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for sub in shape.shapes:
                    shapes_html.append(render_shape(sub))
            else:
                shapes_html.append(render_shape(shape))
        slides_html.append(
            f'<section class="pslide" style="width:{sw}px;height:{sh}px;">'
            f'<div class="pslide-idx">{i:02d} / {len(prs.slides):02d}</div>'
            f'{"".join(shapes_html)}</section>'
        )

    doc = f'''<meta charset="UTF-8">
<title>{html.escape(pptx_path)} — auto pptx→html</title>
<style>
*{{box-sizing:border-box;}}
body{{background:#0A0A0B;color:#E8E6DE;font-family:sans-serif;padding:40px 20px;margin:0;}}
.deck{{max-width:{sw}px;margin:0 auto;display:flex;flex-direction:column;gap:40px;}}
.pslide{{position:relative;background:#18181D;border:1px solid rgba(201,168,76,.3);overflow:hidden;flex-shrink:0;}}
.pslide-idx{{position:absolute;top:4px;right:8px;font-size:11px;color:#9A9890;z-index:99;}}
.pshape{{overflow:hidden;}}
.ptext{{font-size:14px;line-height:1.3;}}
.ptable, .ole-table{{border-collapse:collapse;font-size:11px;width:100%;}}
.ptable td, .ole-table td, .ole-table th{{border:1px solid rgba(255,255,255,.15);padding:3px 6px;}}
.ole-table th{{background:#C9A84C;color:#0A0A0B;}}
</style>
<div class="deck">
{"".join(slides_html)}
</div>
'''
    with open(out_path, 'w', encoding='utf-8') as f:
        f.write(doc)
    print(f'{len(prs.slides)} slides -> {out_path}')


if __name__ == '__main__':
    if len(sys.argv) != 3:
        print('Usage: pptx_to_html.py input.pptx output.html')
        sys.exit(1)
    convert(sys.argv[1], sys.argv[2])
